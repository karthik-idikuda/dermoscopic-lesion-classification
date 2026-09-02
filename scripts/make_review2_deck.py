#!/usr/bin/env python3
"""Fill the faculty Review-2 PPTX template with this project's content.

    python scripts/make_review2_deck.py

Reads the supplied template, replaces the placeholder text in each of the 13
slides, fills the four tables, embeds the generated diagrams, and writes

    docs/review/MP1-UDP-Review2-Dermoscopic-Lesion-Analysis.pptx

The template's own theme, fonts, layouts and decorative pictures are preserved:
only text frames and table cells are rewritten, so the deck still looks like the
one the department issued.

Placeholders that only the student can supply (names, UDP ID, guide, and the
verbatim Review-1 remarks) are written as << ... >> so they are impossible to
miss on screen.
"""

from __future__ import annotations

import copy
from pathlib import Path

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.util import Emu, Pt

ROOT = Path(__file__).resolve().parent.parent
DOCS = ROOT / "docs"
FIGURES = DOCS / "figures"
REVIEW = DOCS / "review"
TEMPLATE = REVIEW / "review2-template.pptx"
OUT = REVIEW / "MP1-UDP-Review2-Dermoscopic-Lesion-Analysis.pptx"

INK = RGBColor(0x10, 0x18, 0x23)
MUTED = RGBColor(0x4D, 0x5B, 0x6E)
TEAL = RGBColor(0x0D, 0x94, 0x88)
ROSE = RGBColor(0xBE, 0x12, 0x3C)
AMBER = RGBColor(0xB4, 0x53, 0x09)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)


# --------------------------------------------------------------------------- #
# Text helpers
# --------------------------------------------------------------------------- #


def set_text(shape, lines, *, size=13, color=INK, bold_first=False, bullet=False,
             space_after=4, line_spacing=None):
    """Replace a shape's text with ``lines``, one paragraph per entry.

    ``lines`` entries may be a plain string or a ``(text, indent_level)`` pair.
    The first paragraph of the existing frame is reused so the template's own
    paragraph properties (alignment, bullet glyphs) survive where possible.
    """
    tf = shape.text_frame
    tf.word_wrap = True

    # Drop every paragraph except the first, which we reuse.
    for p in list(tf.paragraphs)[1:]:
        p._p.getparent().remove(p._p)
    first = tf.paragraphs[0]
    for r in list(first.runs):
        r._r.getparent().remove(r._r)

    for i, entry in enumerate(lines):
        text, level = entry if isinstance(entry, tuple) else (entry, 0)
        para = first if i == 0 else tf.add_paragraph()
        para.level = level
        run = para.add_run()
        run.text = text
        f = run.font
        f.size = Pt(size - 1.0 * level)
        f.color.rgb = color
        f.bold = bold_first and i == 0
        if not bullet:
            # Suppress the inherited bullet glyph for clean prose blocks.
            pPr = para._p.get_or_add_pPr()
            for tag in ("a:buChar", "a:buAutoNum"):
                for el in pPr.findall(f"{{http://schemas.openxmlformats.org/drawingml/2006/main}}{tag.split(':')[1]}"):
                    pPr.remove(el)
        para.space_after = Pt(space_after)
        if line_spacing:
            para.line_spacing = line_spacing
    return tf


def style_cell(cell, text, *, size=9.5, bold=False, color=INK, fill=None):
    cell.text = ""
    tf = cell.text_frame
    tf.word_wrap = True
    para = tf.paragraphs[0]
    run = para.add_run()
    run.text = text
    run.font.size = Pt(size)
    run.font.bold = bold
    run.font.color.rgb = color
    cell.margin_left = Emu(45720)
    cell.margin_right = Emu(45720)
    cell.margin_top = Emu(18288)
    cell.margin_bottom = Emu(18288)
    if fill is not None:
        cell.fill.solid()
        cell.fill.fore_color.rgb = fill


def fit_table(table, rows_needed):
    """Grow or shrink a template table to exactly ``rows_needed`` body rows.

    ``table.rows`` does not support negative indexing, hence the explicit
    ``len() - 1`` lookups.
    """
    body = len(table.rows) - 1
    while body < rows_needed:
        last = table.rows[len(table.rows) - 1]._tr
        table._tbl.append(copy.deepcopy(last))
        body += 1
    while body > rows_needed:
        table._tbl.remove(table.rows[len(table.rows) - 1]._tr)
        body -= 1


def fill_table(shape, header, rows, *, header_size=9.5, body_size=9.0,
               header_fill=TEAL, col_widths=None):
    table = shape.table
    fit_table(table, len(rows))

    if col_widths:
        total = sum(col_widths)
        available = sum(c.width for c in table.columns)
        for col, frac in zip(table.columns, col_widths):
            col.width = Emu(int(available * frac / total))

    for c, label in enumerate(header):
        style_cell(table.cell(0, c), label, size=header_size, bold=True,
                   color=WHITE, fill=header_fill)
    for r, row in enumerate(rows, start=1):
        for c, value in enumerate(row):
            style_cell(table.cell(r, c), value, size=body_size)
    return table


def add_picture(slide, image: Path, *, left, top, width):
    if not image.exists():
        print(f"  ! missing {image.name}, skipped")
        return None
    return slide.shapes.add_picture(str(image), left, top, width=width)


def shrink(shape, *, left=None, top=None, width=None, height=None):
    if left is not None:
        shape.left = left
    if top is not None:
        shape.top = top
    if width is not None:
        shape.width = width
    if height is not None:
        shape.height = height


def table_shape(slide):
    for sh in slide.shapes:
        if sh.has_table:
            return sh
    return None


def text_shapes(slide):
    """Every shape on the slide that can hold text, in document order.

    This template was authored in LibreOffice and gives *every* placeholder
    ``idx == 0``, so placeholder index cannot distinguish title from body.
    Document order can: the title is always first.
    """
    return [sh for sh in slide.shapes if sh.has_text_frame]


def title_of(slide):
    shapes = text_shapes(slide)
    return shapes[0] if shapes else None


def body_placeholder(slide):
    """The first text-bearing shape after the title."""
    shapes = text_shapes(slide)
    return shapes[1] if len(shapes) > 1 else None


# --------------------------------------------------------------------------- #
# Slide content
# --------------------------------------------------------------------------- #

EMU_IN = 914400


def build() -> None:
    prs = Presentation(str(TEMPLATE))
    S = prs.slides
    SW, SH = prs.slide_width, prs.slide_height

    # ---------------------------------------------------------------- slide 1
    s = S[0]
    for sh in s.shapes:
        if sh.has_text_frame and "Presented By" in sh.text_frame.text:
            set_text(sh, [
                "Presented By  —  << your name >>, << team member >>",
                "Project Name :  Explainable Dermoscopic Lesion Analysis",
                "                        EfficientNet-B3 · Grad-CAM · ABCD morphometry · severity grading",
                "Mentor / Guide Name :  << guide name >>",
                "UDP ID  -  U<< id >>",
            ], size=13, color=INK, space_after=7)
        elif sh.has_text_frame and "Review - 2" in sh.text_frame.text:
            set_text(sh, ["Major Project – II", "Review - 2"], size=15, color=WHITE,
                     bold_first=True, space_after=3)

    # ---------------------------------------------------------------- slide 2
    s = S[1]
    body = body_placeholder(s)
    set_text(body, [
        "Progress Since Review-1",
        "Project Objectives",
        "System Design / Architecture",
        "Module-wise Progress",
        "Implementation Progress",
        "AI Tools Used",
        "Challenges and Solutions",
        "Partial Output",
        "Plan for Review-3",
        "References",
    ], size=14, bullet=True, space_after=6)

    # ---------------------------------------------------------------- slide 3
    s = S[2]
    body = body_placeholder(s)
    set_text(body, [
        "Work completed before Review-1",
        ("HAM10000 studied: 10,015 images, 7 diagnostic classes, 67% melanocytic nevi", 1),
        ("Exploratory analysis + SVM baseline (HOG + colour histogram)", 1),
        ("EfficientNet-B3 training notebook, reported 80.17% accuracy", 1),
        "Work completed after Review-1",
        ("Notebooks refactored into an installable package: 18 modules, 5,356 lines in src/derm", 1),
        ("FastAPI backend with 12 endpoints + browser workstation UI (3,353 lines)", 1),
        ("Explainability: Grad-CAM / Grad-CAM++ scored against the lesion mask", 1),
        ("Quantitative Stolz ABCD morphometry as an independent, non-neural evidence stream", 1),
        ("Uncertainty: TTA, MC dropout, predictive entropy, BALD", 1),
        ("Severity grading (0-100) with one-directional safety overrides", 1),
        ("Rule-based clinical narrative + PDF report, SQLite case store, change tracking", 1),
        ("Reproducible data-leakage audit; 149 automated tests, all passing", 1),
    ], size=11, bullet=True, space_after=2)
    shrink(body, top=Emu(int(1.15 * EMU_IN)), height=Emu(int(2.55 * EMU_IN)))

    tbl = table_shape(s)
    if tbl:
        shrink(tbl, left=Emu(int(0.45 * EMU_IN)), top=Emu(int(3.78 * EMU_IN)),
               width=Emu(int(9.1 * EMU_IN)))
        fill_table(tbl,
                   ["Review-1 Remark", "Action Taken", "Status"],
                   [
                       ["1.  << paste your actual Review-1 remark >>  e.g. “justify the "
                        "reported accuracy”",
                        "Ran a reproducible leakage audit: the notebook's image-wise split "
                        "leaks 36.13% of test images. Switched to lesion-grouped splitting "
                        "(0% leakage) and tagged all notebook figures as self-reported.",
                        "Completed"],
                       ["2.  << paste remark >>  e.g. “model must be explainable and have a "
                        "working front end”",
                        "Added Grad-CAM/Grad-CAM++ with an attention-alignment score, ABCD "
                        "morphometry, and a full FastAPI + browser application replacing the "
                        "notebook-only workflow.",
                        "Completed"],
                   ],
                   header_size=9, body_size=8, col_widths=[0.26, 0.60, 0.14])

    # ---------------------------------------------------------------- slide 4
    s = S[3]
    body = body_placeholder(s)
    set_text(body, [
        "Original objectives",
        ("Classify dermoscopic images into the 7 HAM10000 categories with EfficientNet-B3", 1),
        ("Benchmark against a classical SVM baseline", 1),
        ("Produce a Grad-CAM visual explanation for each prediction", 1),
        ("Derive an automated severity indication from the prediction", 1),
        "Updated objectives (added after Review-1)",
        ("Quantitative Stolz ABCD morphometry — evidence that does not depend on learned weights", 1),
        ("Uncertainty quantification so unsure cases are escalated, not silently guessed", 1),
        ("Longitudinal change tracking — the “E” of ABCDE, which a single image cannot show", 1),
        ("Deterministic clinical narrative + PDF report for handover", 1),
        ("Reproducible data-integrity audit and a checkpoint verifier", 1),
        "Scope modifications",
        ("Splitting changed from image-wise to lesion-grouped, to remove near-duplicate leakage", 1),
        ("Training moved off local hardware: an 8 GB M2 drove swap to 13.9 GB (documented, not hidden)", 1),
        ("Narrative kept rule-based rather than an LLM, so every sentence traces to a measurement", 1),
    ], size=10.5, bullet=True, space_after=1.5)

    # ---------------------------------------------------------------- slide 5
    s = S[4]
    body = body_placeholder(s)
    if body:
        body._element.getparent().remove(body._element)
    for sh in list(s.shapes):
        if sh.shape_type == 13 and not sh.has_text_frame:  # decorative picture
            sh._element.getparent().remove(sh._element)
    add_picture(s, FIGURES / "architecture.png",
                left=Emu(int(0.30 * EMU_IN)), top=Emu(int(1.02 * EMU_IN)),
                width=Emu(int(9.40 * EMU_IN)))

    # a second architecture slide for the workflow + database design
    def clone_slide_after(index, title_text):
        src = S[index]
        new = prs.slides.add_slide(src.slide_layout)
        # Keep only the title; the diagram fills the rest of the slide.
        for sh in text_shapes(new)[1:]:
            sh._element.getparent().remove(sh._element)
        t = title_of(new)
        if t is not None:
            set_text(t, [title_text], size=20, bold_first=True)
        # move it directly after `index`
        xml_slides = prs.slides._sldIdLst
        slides = list(xml_slides)
        xml_slides.remove(slides[-1])
        xml_slides.insert(index + 1, slides[-1])
        return new

    wf = clone_slide_after(4, "Workflow Diagram")
    add_picture(wf, FIGURES / "workflow.png",
                left=Emu(int(0.30 * EMU_IN)), top=Emu(int(1.02 * EMU_IN)),
                width=Emu(int(9.40 * EMU_IN)))

    db = clone_slide_after(5, "Database Design")
    add_picture(db, FIGURES / "er_diagram.png",
                left=Emu(int(0.75 * EMU_IN)), top=Emu(int(1.02 * EMU_IN)),
                width=Emu(int(8.50 * EMU_IN)))

    # indices shift by 2 from here on
    off = 2

    # ---------------------------------------------------------------- slide 6
    s = S[5 + off]
    tbl = table_shape(s)
    if tbl:
        shrink(tbl, left=Emu(int(0.35 * EMU_IN)), top=Emu(int(1.05 * EMU_IN)),
               width=Emu(int(9.3 * EMU_IN)))
        fill_table(tbl,
                   ["Module", "Description", "Status"],
                   [
                       ["1. Preprocessing", "Hair detection + Telea inpainting, Shades-of-Gray "
                        "colour constancy, lens-vignette crop, EXIF handling", "Complete"],
                       ["2. Quality & OOD gate", "Focus (variance of Laplacian), exposure, contrast, "
                        "glare, and a skin-chromaticity test that rejects non-skin images", "Complete"],
                       ["3. Segmentation", "Lesion-enhanced Otsu, morphological cleanup, "
                        "centrality-weighted component, ellipse fallback with reduced confidence", "Complete"],
                       ["4. ABCD morphometry", "Asymmetry about principal axes, 8-sector border "
                        "irregularity, six-colour counting, structure approximation → TDS", "Complete"],
                       ["5. Classification", "EfficientNet-B3 (timm), checkpoint bundle with graceful "
                        "untrained fallback and explicit weights_status", "Complete"],
                       ["6. Explainability", "Grad-CAM and Grad-CAM++ with hook cleanup, scored for "
                        "attention alignment against the lesion mask", "Complete"],
                       ["7. Uncertainty", "Dihedral TTA, MC dropout, predictive entropy, BALD mutual "
                        "information, verdict classification", "Complete"],
                       ["8. Severity grading", "Composite 0-100 score with one-directional melanoma / "
                        "TDS / low-confidence safety overrides", "Complete"],
                       ["9. Reporting", "Deterministic rule-based narrative and ReportLab PDF export", "Complete"],
                       ["10. Change tracking", "Two-capture longitudinal comparison with optional "
                        "millimetre calibration", "Complete"],
                       ["11. Persistence", "SQLite case store: history, filters, aggregate statistics, "
                        "annotations", "Complete"],
                       ["12. Backend", "FastAPI service, 12 endpoints, Pydantic validation, optional "
                        "API-key auth", "Complete"],
                       ["13. Frontend", "Vanilla-JS clinical workstation: studio, batch triage, "
                        "tracking, history, metrics", "Complete"],
                       ["14. Validation", "149 automated tests on synthetic fixtures + smoke test + "
                        "checkpoint verifier", "Complete"],
                       ["15. Data integrity", "Reproducible leakage audit written to "
                        "docs/split_audit.json", "Complete"],
                       ["16. Dataset acquisition", "All 10,015 HAM10000 images downloaded from "
                        "Harvard Dataverse and downscaled to 195 MB on disk", "Complete"],
                       ["17. Trained weights", "Classifier head fitted on frozen ImageNet "
                        "EfficientNet-B3 features (linear probe); passes all 21 checkpoint "
                        "verification checks", "Complete"],
                       ["18. Evaluation", "Measured on the leak-free test split: 53.7% balanced "
                        "accuracy, 0.379 macro-F1, ROC-AUC 0.838, ECE 0.096", "Complete"],
                       ["19. End-to-end fine-tune", "Needs a CUDA GPU; the linear probe is the "
                        "CPU-feasible substitute", "Review-3"],
                       ["20. Deployment", "Containerisation, auth hardening and hosting", "Review-3"],
                   ],
                   header_size=9, body_size=7.2, col_widths=[0.20, 0.66, 0.14])

    # ---------------------------------------------------------------- slide 7
    s = S[6 + off]
    body = body_placeholder(s)
    set_text(body, [
        "Technologies: Python 3.11 · PyTorch 2 + timm · OpenCV · scikit-image · scikit-learn · "
        "FastAPI + Uvicorn · Pydantic v2 · ReportLab · SQLite · vanilla HTML/CSS/ES2020 (no build step)",
        "Overall implementation: ~9,000 lines of application code and ~1,300 lines of tests.",
    ], size=10, color=MUTED, space_after=3)
    shrink(body, top=Emu(int(1.02 * EMU_IN)), height=Emu(int(0.85 * EMU_IN)))

    tbl = table_shape(s)
    if tbl:
        shrink(tbl, left=Emu(int(1.55 * EMU_IN)), top=Emu(int(2.00 * EMU_IN)),
               width=Emu(int(6.9 * EMU_IN)))
        fill_table(tbl,
                   ["Module", "Completion"],
                   [
                       ["Image processing pipeline (preprocess, quality, segmentation, ABCD)", "100%"],
                       ["Explainability (Grad-CAM / Grad-CAM++, attention alignment)", "100%"],
                       ["Uncertainty + severity grading", "100%"],
                       ["Reporting, change tracking, case store", "100%"],
                       ["FastAPI backend (12 endpoints)", "100%"],
                       ["Frontend workstation UI", "100%"],
                       ["Automated tests (149 passing) + leakage audit", "100%"],
                       ["Dataset acquisition — all 10,015 images", "100%"],
                       ["Trained weights (linear probe) + evaluation", "100%"],
                       ["End-to-end fine-tune (needs a CUDA GPU)", "40%"],
                       ["Deployment & packaging", "35%"],
                   ],
                   header_size=10, body_size=9, col_widths=[0.78, 0.22])

    # ---------------------------------------------------------------- slide 8
    s = S[7 + off]
    tbl = table_shape(s)
    if tbl:
        shrink(tbl, left=Emu(int(0.30 * EMU_IN)), top=Emu(int(1.00 * EMU_IN)),
               width=Emu(int(9.4 * EMU_IN)))
        fill_table(tbl,
                   ["AI Tool", "Purpose", "Sample Prompt", "Outcome"],
                   [
                       ["Kiro\n(Claude Opus 5)\nAI IDE agent",
                        "Primary development assistant: refactoring the notebooks into a "
                        "package, writing the FastAPI service and the frontend, and "
                        "generating the test suite.",
                        "“Analyse this HAM10000 notebook project in detail and convert it "
                        "into a working frontend + backend application with Grad-CAM, ABCD "
                        "morphometry and severity grading; all features must actually work.”",
                        "18-module src/derm package, 12-endpoint API and full UI. Verified by "
                        "149 passing tests and an end-to-end smoke test rather than by "
                        "reading the code alone."],
                       ["Kiro\n(agentic debugging)",
                        "Root-causing defects that unit tests surfaced in the inference "
                        "pipeline.",
                        "“MC dropout returns a standard deviation of ~1e-6, so the "
                        "uncertainty estimate is inert. Find out why and fix it.”",
                        "Found that timm's EfficientNet applies dropout functionally, so "
                        "model.train() alone did not make passes stochastic. Fix confirmed "
                        "by an assertion that MC variance is non-zero."],
                       ["Kiro\n(data audit)",
                        "Independently checking the accuracy claim carried over from "
                        "Review-1.",
                        "“HAM10000 has repeated photographs of the same lesion. Quantify how "
                        "much of the notebook's test set leaks from training, per class.”",
                        "Produced scripts/audit_leakage.py and docs/split_audit.json: 36.13% "
                        "of test images leak, 63.5% for melanoma. Reproducible from the "
                        "metadata CSV alone, so the finding is checkable, not asserted."],
                       ["Kiro\n(working within a\nhardware limit)",
                        "Obtaining real trained weights without a GPU, after full "
                        "fine-tuning exhausted the machine's memory.",
                        "“Don't train images on my Mac. Download the dataset and make the "
                        "whole project work end to end.”",
                        "Split transfer learning in two: forward-only feature extraction "
                        "(bounded RAM, 472 MB peak) then a linear head fit in seconds. "
                        "Balanced accuracy 14.3% → 53.7%. Verified by 21 checkpoint checks "
                        "including per-class recall on real images, which is what proves the "
                        "class order is not permuted."],
                       ["ChatGPT / Claude\n(chat)",
                        "Clarifying the clinical basis of the Stolz ABCD rule and the "
                        "Grad-CAM++ weighting before implementing them.",
                        "“Explain how the Stolz ABCD total dermoscopy score is weighted, and "
                        "how Grad-CAM++ differs from Grad-CAM.”",
                        "Implementation cross-checked against the original papers; the TDS "
                        "weights and the >5.45 threshold were taken from the literature, not "
                        "from the model's suggestion."],
                       ["GitHub Copilot",
                        "Inline completion for repetitive code: dataclass fields, argparse "
                        "wiring, docstrings.",
                        "(inline completion, no discrete prompt)",
                        "Reduced boilerplate typing. All completions reviewed before commit; "
                        "no completion was accepted into numerical logic unchecked."],
                       ["<< remove any row your team did not actually use >>", "", "", ""],
                   ],
                   header_size=8.5, body_size=6.4,
                   col_widths=[0.13, 0.22, 0.34, 0.31])

    # ---------------------------------------------------------------- slide 9
    s = S[8 + off]
    tbl = table_shape(s)
    if tbl:
        shrink(tbl, left=Emu(int(0.35 * EMU_IN)), top=Emu(int(1.02 * EMU_IN)),
               width=Emu(int(9.3 * EMU_IN)))
        fill_table(tbl,
                   ["Challenge", "Solution Implemented"],
                   [
                       ["Reported 80.17% accuracy could not be trusted: HAM10000 holds 10,015 "
                        "images of only 7,470 lesions, so an image-wise split puts "
                        "near-duplicate photographs of the same lesion in train and test.",
                        "Wrote a metadata-only audit (scripts/audit_leakage.py). Measured "
                        "36.13% test leakage overall, 70.1% for BCC and 63.5% for melanoma. "
                        "Switched to lesion-grouped splitting, which audits at 0.00%. All "
                        "unverified figures are now tagged self_reported in the UI."],
                       ["Severe class imbalance — 67% of the dataset is melanocytic nevi, so "
                        "plain accuracy rewards a model that ignores melanoma.",
                        "Checkpoint selection on macro-F1 instead of accuracy, class-weighted "
                        "loss, and balanced accuracy plus per-class recall reported alongside "
                        "a melanoma safety-net catch rate."],
                       ["Training EfficientNet-B3 locally was not viable: on an 8 GB M2 the MPS "
                        "backend drove swap to 13.9 GB and cut free disk from 9.2 GB to 2.9 GB "
                        "in ninety seconds.",
                        "Diagnosed as a hardware limit, not a code defect, and documented. "
                        "Added scripts/bench.py to measure throughput before committing, and a "
                        "disk-frugal downloader that streams one archive at a time and "
                        "downscales on extraction (2.9 GB → ~350 MB)."],
                       ["A missing or mismatched checkpoint could silently yield confident, "
                        "plausible, wrong diagnoses — the most dangerous failure mode here.",
                        "weights_status is surfaced through the API and the UI, and the app "
                        "refuses to present untrained output as meaningful. "
                        "scripts/verify_checkpoint.py additionally rejects wrong class counts, "
                        "permuted class order, and uniform (untrained) outputs."],
                       ["MC dropout produced ~1e-6 variance, making the uncertainty estimate "
                        "meaningless; Grad-CAM also failed inside torch.inference_mode().",
                        "timm applies dropout functionally, so train() alone was insufficient — "
                        "stochasticity is now activated explicitly. Grad-CAM escapes the outer "
                        "inference-mode context so gradients can flow. Both are covered by "
                        "regression tests."],
                       ["Hair, glare and the black lens barrel were being read as pigment, "
                        "corrupting the ABCD measurements.",
                        "Directional black-hat hair detection with Telea inpainting plus "
                        "vignette cropping. The restored frame drives geometry only — the "
                        "classifier still receives the un-restored image, avoiding a silent "
                        "train/serve skew."],
                   ],
                   header_size=9, body_size=7.0, col_widths=[0.42, 0.58])

    # --------------------------------------------------------------- slide 10
    s = S[9 + off]
    body = body_placeholder(s)
    set_text(body, [
        "Working end-to-end today, with trained weights loaded",
        ("Demo flow: upload a dermoscopic image → quality gate → segmentation overlay → "
         "ABCD rings and TDS → ranked class probabilities → Grad-CAM split-view against the "
         "original → severity tier with recommended action → clinical narrative → PDF download.", 1),
        ("Batch triage: many images scored and returned in severity order, so review time "
         "goes where the risk is.", 1),
        ("Track change: two captures of one lesion compared, with millimetre measurement "
         "when a field-of-view width is supplied.", 1),
        ("Case history: persisted cases with tier/class filters, aggregate statistics and "
         "annotation.", 1),
        ("Model metrics: measured vs self-reported provenance, plus the leakage audit "
         "rendered in-app.", 1),
        "Measured results on the leak-free test split (1,523 images, 1,120 lesions)",
        ("Balanced accuracy 53.7% · accuracy 55.0% · macro-F1 0.379 · ROC-AUC 0.838 · "
         "ECE 0.096 after temperature scaling", 1),
        ("Melanoma recall 55.5%; the severity safety net escalates 59.2% of true melanomas "
         "to HIGH/CRITICAL, at a 15.3% benign over-referral cost.", 1),
        ("Chance is 14.3% and a trivial always-nevus model scores 14.3% balanced accuracy, "
         "so the head has genuinely learned.", 1),
        ("Weights come from a LINEAR PROBE: the classifier head was fitted on frozen "
         "ImageNet EfficientNet-B3 features. The backbone was not fine-tuned, because that "
         "needs a GPU. Reported as a probe, never as a fine-tuned model.", 1),
        ("149/149 automated tests passing; smoke test green across all 12 endpoints; "
         "checkpoint passes all 21 verification checks including class-order confirmation "
         "on real images.", 1),
        ("Leakage audit reproducible from metadata: 36.13% image-wise vs 0.00% lesion-grouped.", 1),
        ("Full-pipeline CPU inference ≈ 2.9 s per image on an Apple M2.", 1),
        "<< paste UI screenshots here: analysis studio, Grad-CAM split view, batch triage, PDF report >>",
    ], size=9.5, bullet=True, space_after=1.5)

    # --------------------------------------------------------------- slide 11
    s = S[10 + off]
    body = body_placeholder(s)
    set_text(body, [
        "Remaining module implementation",
        ("Fine-tune the full backbone end to end on Kaggle / Colab / a CUDA GPU — the linear "
         "probe's 53.7% balanced accuracy is the CPU ceiling, not the method's", 1),
        ("Expect the biggest gain on the rare classes: df precision is 0.065 and akiec 0.20, "
         "because a frozen ImageNet backbone has no dermoscopy-specific features", 1),
        ("Re-verify with scripts/verify_checkpoint.py --images before trusting the result", 1),
        "Testing",
        ("Extend the 149-test suite with trained-model integration tests and per-class recall floors", 1),
        ("Re-run the melanoma safety-net audit against real weights", 1),
        "Integration & documentation",
        ("Regenerate docs/evaluation.json, confusion matrix and calibration figures from the real run", 1),
        ("Replace every self_reported figure with a measured one, or state plainly that it is unverified", 1),
        ("Final project report and user guide", 1),
        "Performance optimisation",
        ("ONNX export / quantisation to cut the 2.9 s CPU inference time", 1),
        ("Response caching for repeated figure and metrics requests", 1),
        "Final deployment",
        ("Containerise, enable DERM_API_KEY, serve over HTTPS, add structured logging", 1),
    ], size=10.5, bullet=True, space_after=1.5)

    # --------------------------------------------------------------- slide 12
    s = S[11 + off]
    body = body_placeholder(s)
    set_text(body, [
        "1.  Tschandl P., Rosendahl C., Kittler H. The HAM10000 dataset, a large collection of "
        "multi-source dermatoscopic images of common pigmented skin lesions. Scientific Data 5, "
        "180161 (2018). doi:10.1038/sdata.2018.161",
        "2.  Harvard Dataverse, HAM10000 distribution. doi:10.7910/DVN/DBW86T (CC BY-NC 4.0).",
        "3.  Tan M., Le Q. EfficientNet: Rethinking Model Scaling for Convolutional Neural "
        "Networks. ICML 2019.",
        "4.  Selvaraju R. R. et al. Grad-CAM: Visual Explanations from Deep Networks via "
        "Gradient-based Localization. ICCV 2017.",
        "5.  Chattopadhyay A. et al. Grad-CAM++: Improved Visual Explanations for Deep "
        "Convolutional Networks. WACV 2018.",
        "6.  Gal Y., Ghahramani Z. Dropout as a Bayesian Approximation: Representing Model "
        "Uncertainty in Deep Learning. ICML 2016.",
        "7.  Houlsby N. et al. Bayesian Active Learning for Classification and Preference "
        "Learning. arXiv:1112.5745 (2011).",
        "8.  Guo C. et al. On Calibration of Modern Neural Networks. ICML 2017.",
        "9.  Stolz W. et al. ABCD rule of dermatoscopy: a new practical method for early "
        "recognition of malignant melanoma. European Journal of Dermatology, 1994.",
        "10. Telea A. An Image Inpainting Technique Based on the Fast Marching Method. "
        "Journal of Graphics Tools, 2004.",
        "11. Finlayson G., Trezzi E. Shades of Gray and Colour Constancy. Color Imaging "
        "Conference, 2004.",
        "12. Wightman R. PyTorch Image Models (timm). github.com/huggingface/pytorch-image-models",
    ], size=9.5, space_after=3.5)

    # --------------------------------------------------------------- slide 13
    s = S[12 + off]
    phs = [sh for sh in s.shapes if sh.has_text_frame]
    if phs:
        set_text(phs[0], ["Thank You"], size=34, bold_first=True, color=INK)
    if len(phs) > 1:
        set_text(phs[1], ["Questions & Demo"], size=16, color=MUTED)

    DOCS.mkdir(parents=True, exist_ok=True)
    prs.save(str(OUT))
    print(f"Wrote {OUT.relative_to(ROOT)}  ({len(prs.slides._sldIdLst)} slides)")


if __name__ == "__main__":
    build()
